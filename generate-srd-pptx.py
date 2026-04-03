#!/usr/bin/env python3
"""
Generate SRD (Strategic Readiness Discussion) Analysis PowerPoint Presentation.

Generates a branded Dale Carnegie Queensland PPTX from SRD data stored in srd-data.json.
New companies can be added as additional objects in the JSON array.

Usage:
    python3 generate-srd-pptx.py                          # Uses first company in JSON
    python3 generate-srd-pptx.py "Golden Cockerel Chicken" # Specify company by name
"""

import os
import io
import json
import math
import sys
import statistics
import tempfile
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

# ─── Brand Colours ───────────────────────────────────────────────────────────
BLACK   = RGBColor(0x00, 0x00, 0x00)
YELLOW  = RGBColor(0xFF, 0xC7, 0x08)
GREEN   = RGBColor(0x84, 0xC4, 0x4C)
BLUE    = RGBColor(0x00, 0x90, 0xCF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0x99, 0x99, 0x99)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)

# Logos
LOGO_WHITE = os.path.expanduser("~/Dropbox/2026/Marketing/Logos & Badges/White DC Logo.png")
LOGO_BLACK = os.path.expanduser("~/Dropbox/2026/Marketing/Logos & Badges/Black DC Logo No Background.png")
CLIENT_LOGOS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "client-logos")
COVER_BG = os.path.expanduser("~/Desktop/slide1_img_0.png")

# Header bar colour (DC Green)
HEADER_BG = GREEN

# Score colour coding
AMBER = RGBColor(0xFF, 0xA5, 0x00)
RED   = RGBColor(0xFF, 0x63, 0x47)

def score_colour(score):
    """Return RGBColor based on score range."""
    if score >= 9:
        return GREEN
    elif score >= 7:
        return YELLOW
    elif score >= 5:
        return AMBER
    else:
        return RED

def score_text_colour(score):
    """White text on most cells, black on yellow."""
    if 7 <= score <= 8:
        return BLACK
    return WHITE

def avg_score_colour(avg):
    """Score colour for float averages."""
    if avg >= 9:
        return GREEN
    elif avg >= 7:
        return YELLOW
    elif avg >= 5:
        return AMBER
    else:
        return RED


# ─── Profile Names ─────────────────────────────────────────────────────────

PROFILES = ["Leadership DNA", "UYLP", "Leadership for Results"]


# ─── Data Loading ──────────────────────────────────────────────────────────

def load_data(company_name=None):
    """Load company data from srd-data.json, converting dicts to OrderedDicts."""
    json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "srd-data.json")
    with open(json_path) as f:
        all_data = json.load(f)
    companies = all_data["companies"]
    if company_name:
        for c in companies:
            if c["company"].lower() == company_name.lower():
                return _convert_to_ordered(c)
        print(f"Company '{company_name}' not found. Available: {[c['company'] for c in companies]}")
        sys.exit(1)
    return _convert_to_ordered(companies[0])


def _convert_to_ordered(company_data):
    """Convert profile score dicts to OrderedDicts for consistent ordering."""
    for participant in company_data["participants"]:
        for profile in PROFILES:
            if profile in participant:
                participant[profile] = OrderedDict(participant[profile])
    return company_data


# ─── Helper Functions ────────────────────────────────────────────────────────

def compute_team_averages(data, profile_name):
    """Compute team average for each competency in a profile."""
    competencies = list(data["participants"][0][profile_name].keys())
    averages = OrderedDict()
    for comp in competencies:
        scores = [p[profile_name][comp] for p in data["participants"]]
        averages[comp] = round(statistics.mean(scores), 1)
    return averages


def compute_participant_profile_avg(participant, profile_name):
    """Average score for one participant across one profile."""
    scores = list(participant[profile_name].values())
    return round(statistics.mean(scores), 1)


def compute_overall_avg(participant):
    """Average score across all three profiles."""
    all_scores = []
    for profile in PROFILES:
        all_scores.extend(participant[profile].values())
    return round(statistics.mean(all_scores), 1)


def get_all_competency_averages(data):
    """Get average score for every competency across all profiles, return sorted list."""
    results = []
    for profile in PROFILES:
        avgs = compute_team_averages(data, profile)
        for comp, avg in avgs.items():
            results.append((comp, avg, profile))
    return results


def set_cell_text(cell, text, font_size=Pt(8), bold=False, colour=BLACK, alignment=PP_ALIGN.CENTER):
    """Set text in a table cell with formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Arial"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_shape_rect(slide, left, top, width, height, fill_colour):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 colour=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = colour
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_slide_header(slide, title_text, font_size=Pt(18)):
    """Add a green header bar with white logo and title — generous spacing."""
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 title_text, font_size=font_size, colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))


def add_slide_footer(slide):
    """Add a green accent line at the bottom of a slide."""
    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


def add_slide_number(slide, number):
    """Add slide number to bottom right corner."""
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
                 str(number), font_size=Pt(10), colour=GREY, bold=False, alignment=PP_ALIGN.RIGHT)


def first_name(full_name):
    """Get short display name."""
    parts = full_name.split()
    if len(parts) >= 2:
        return parts[0][0] + ". " + parts[-1]
    return full_name


# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    """Slide 1: Clean green cover with partnership logos equally sized at the top."""
    print("  Building Slide 1: Title...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Solid green background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    # ── Logos: large, spread wide on each side ──
    logo_height = Inches(1.8)
    logo_y = Inches(0.5)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(1.0), logo_y, height=logo_height)
    client_logo = os.path.join(os.path.dirname(os.path.abspath(__file__)), data.get("logo", "")) if data.get("logo") else None
    if client_logo and os.path.exists(client_logo):
        slide.shapes.add_picture(client_logo, Inches(10.5), logo_y, height=logo_height)

    # Main title — centred, pushed down for breathing room
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
                 "Strategic Readiness\nDiscussion Analysis",
                 font_size=Pt(40), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Company name
    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.8),
                 data["company"],
                 font_size=Pt(28), colour=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.4),
                 "Prepared by Dale Carnegie Queensland",
                 font_size=Pt(13), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)


def build_team_overview_slide(prs, data):
    """Team overview with strengths and development areas."""
    print("  Building Slide: Team Overview...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar — taller, more breathing room
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 "Team Overview", font_size=Pt(24), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    num_participants = len(data["participants"])
    all_scores = []
    for p in data["participants"]:
        for profile in PROFILES:
            all_scores.extend(p[profile].values())
    team_avg = round(statistics.mean(all_scores), 1)

    # Stats boxes
    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(2.5), Inches(1.2), BLUE)
    add_text_box(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.4),
                 "Participants", font_size=Pt(14), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(2.5), Inches(0.6),
                 str(num_participants), font_size=Pt(36), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(3.3), Inches(1.2), Inches(2.5), Inches(1.2), DARK_GREY)
    add_text_box(slide, Inches(3.3), Inches(1.3), Inches(2.5), Inches(0.4),
                 "Team Average Score", font_size=Pt(14), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.3), Inches(1.7), Inches(2.5), Inches(0.6),
                 f"{team_avg}", font_size=Pt(36), colour=YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

    all_comp = get_all_competency_averages(data)
    sorted_high = sorted(all_comp, key=lambda x: x[1], reverse=True)
    sorted_low = sorted(all_comp, key=lambda x: x[1])

    seen = set()
    top5 = []
    for comp, avg, profile in sorted_high:
        if comp not in seen and len(top5) < 5:
            top5.append((comp, avg, profile))
            seen.add(comp)

    seen = set()
    bottom5 = []
    for comp, avg, profile in sorted_low:
        if comp not in seen and len(bottom5) < 5:
            bottom5.append((comp, avg, profile))
            seen.add(comp)

    # Left column: Top 5 Strengths
    left_x = Inches(0.5)
    add_text_box(slide, left_x, Inches(2.7), Inches(6), Inches(0.5),
                 "Top 5 Strengths", font_size=Pt(18), colour=BLACK, bold=True)

    y = Inches(3.3)
    for comp, avg, profile in top5:
        add_shape_rect(slide, left_x, y + Pt(4), Inches(0.15), Inches(0.25), GREEN)
        add_text_box(slide, left_x + Inches(0.25), y, Inches(4.5), Inches(0.35),
                     f"{comp}  ({avg})", font_size=Pt(13), colour=BLACK, bold=False)
        add_text_box(slide, left_x + Inches(4.8), y, Inches(1.5), Inches(0.35),
                     profile, font_size=Pt(9), colour=GREY, bold=False)
        y += Inches(0.4)

    # Right column: Top 5 Development Areas
    right_x = Inches(7)
    add_text_box(slide, right_x, Inches(2.7), Inches(6), Inches(0.5),
                 "Top 5 Development Areas", font_size=Pt(18), colour=BLACK, bold=True)

    y = Inches(3.3)
    for comp, avg, profile in bottom5:
        add_shape_rect(slide, right_x, y + Pt(4), Inches(0.15), Inches(0.25), RED)
        add_text_box(slide, right_x + Inches(0.25), y, Inches(4.5), Inches(0.35),
                     f"{comp}  ({avg})", font_size=Pt(13), colour=BLACK, bold=False)
        add_text_box(slide, right_x + Inches(4.8), y, Inches(1.5), Inches(0.35),
                     profile, font_size=Pt(9), colour=GREY, bold=False)
        y += Inches(0.4)

    add_shape_rect(slide, Inches(0), Inches(7.2), Inches(13.333), Pt(4), YELLOW)


def build_score_table_slide(prs, data, profile_name):
    """Build a score table slide for a given profile."""
    print(f"  Building Score Table: {profile_name}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar — taller with breathing room
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 f"{profile_name} — Improvement Profile",
                 font_size=Pt(20), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    participants = data["participants"]
    competencies = list(participants[0][profile_name].keys())
    team_avgs = compute_team_averages(data, profile_name)

    num_rows = len(competencies) + 1
    num_cols = len(participants) + 2

    table_left = Inches(0.2)
    table_top = Inches(1.05)
    table_width = Inches(12.9)
    table_height = Inches(6.1)

    tbl_shape = slide.shapes.add_table(num_rows, num_cols, table_left, table_top, table_width, table_height)
    tbl = tbl_shape.table

    comp_col_width = Inches(2.2)
    score_col_width = Inches(1.18)
    avg_col_width = Inches(1.2)

    tbl.columns[0].width = Emu(int(comp_col_width))
    for i in range(1, num_cols - 1):
        tbl.columns[i].width = Emu(int(score_col_width))
    tbl.columns[num_cols - 1].width = Emu(int(avg_col_width))

    # Header row
    set_cell_text(tbl.cell(0, 0), "Competency", Pt(8), True, WHITE, PP_ALIGN.LEFT)
    tbl.cell(0, 0).fill.solid()
    tbl.cell(0, 0).fill.fore_color.rgb = BLACK

    for i, p in enumerate(participants):
        short = first_name(p["name"])
        set_cell_text(tbl.cell(0, i + 1), short, Pt(7), True, WHITE, PP_ALIGN.CENTER)
        tbl.cell(0, i + 1).fill.solid()
        tbl.cell(0, i + 1).fill.fore_color.rgb = BLACK

    set_cell_text(tbl.cell(0, num_cols - 1), "Team Avg", Pt(8), True, WHITE, PP_ALIGN.CENTER)
    tbl.cell(0, num_cols - 1).fill.solid()
    tbl.cell(0, num_cols - 1).fill.fore_color.rgb = BLACK

    # Data rows
    for r, comp in enumerate(competencies):
        row_idx = r + 1
        set_cell_text(tbl.cell(row_idx, 0), comp, Pt(8), False, BLACK, PP_ALIGN.LEFT)
        tbl.cell(row_idx, 0).fill.solid()
        tbl.cell(row_idx, 0).fill.fore_color.rgb = WHITE

        for c, p in enumerate(participants):
            score = p[profile_name][comp]
            col_idx = c + 1
            cell = tbl.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = score_colour(score)
            set_cell_text(cell, str(score), Pt(9), True, score_text_colour(score), PP_ALIGN.CENTER)

        avg = team_avgs[comp]
        avg_cell = tbl.cell(row_idx, num_cols - 1)
        avg_cell.fill.solid()
        avg_cell.fill.fore_color.rgb = avg_score_colour(avg)
        set_cell_text(avg_cell, f"{avg}", Pt(9), True,
                      BLACK if 7 <= avg < 9 else WHITE, PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


def build_bar_chart_slide(prs, data, profile_name):
    """Build a polished horizontal bar chart slide for a profile — print-ready."""
    print(f"  Building Bar Chart: {profile_name}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 f"{profile_name} — Team Averages",
                 font_size=Pt(20), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    team_avgs = compute_team_averages(data, profile_name)
    sorted_avgs = sorted(team_avgs.items(), key=lambda x: x[1], reverse=True)

    num_bars = len(sorted_avgs)
    chart_left = Inches(3.2)
    chart_top = Inches(1.2)
    chart_width = Inches(9.0)
    bar_area_height = Inches(5.8)
    bar_height = bar_area_height / num_bars * 0.65
    bar_gap = bar_area_height / num_bars

    max_score = 10.0

    # Subtle gridlines at 2, 4, 6, 8, 10
    for tick in [2, 4, 6, 8, 10]:
        x = chart_left + int(chart_width * (tick / max_score))
        add_shape_rect(slide, x, chart_top, Pt(1), int(bar_area_height), LIGHT_GREY)
        add_text_box(slide, x - Inches(0.15), chart_top + int(bar_area_height) + Inches(0.05),
                     Inches(0.3), Inches(0.25),
                     str(tick), font_size=Pt(8), colour=GREY, bold=False, alignment=PP_ALIGN.CENTER)

    for i, (comp, avg) in enumerate(sorted_avgs):
        y = chart_top + int(i * bar_gap)
        bar_w = int(chart_width * (avg / max_score))

        # Label — right-aligned to the left of bars
        add_text_box(slide, Inches(0.2), y, Inches(2.9), int(bar_height),
                     comp, font_size=Pt(10), colour=DARK_GREY, bold=False, alignment=PP_ALIGN.RIGHT)

        # Bar with rounded corners (using ROUND_RECTANGLE)
        colour = avg_score_colour(avg)
        bar_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, chart_left, y + Pt(3),
            bar_w, int(bar_height) - Pt(6))
        bar_shape.fill.solid()
        bar_shape.fill.fore_color.rgb = colour
        bar_shape.line.fill.background()

        # Score label at end of bar
        add_text_box(slide, chart_left + bar_w + Inches(0.1), y, Inches(0.6), int(bar_height),
                     f"{avg}", font_size=Pt(11), colour=DARK_GREY, bold=True, alignment=PP_ALIGN.LEFT)

    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


def build_collective_themes_slide(prs, data):
    """Collective themes from strategic questions."""
    print("  Building Slide: Collective Themes...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 "Strategic Readiness — Collective Themes",
                 font_size=Pt(20), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    questions = data["questions"]
    q_keys = ["Q1", "Q2", "Q3", "Q4", "Q5"]

    y = Inches(1.1)
    for qi, (qkey, question) in enumerate(zip(q_keys, questions)):
        all_responses = []
        for p in data["participants"]:
            all_responses.extend(p["strategic"][qkey])

        from collections import Counter
        lower_map = {}
        for r in all_responses:
            key = r.lower().strip()
            if key not in lower_map:
                lower_map[key] = r
        freq = Counter([r.lower().strip() for r in all_responses])
        common = [lower_map[k] for k, v in freq.most_common() if v >= 2]
        singles = [lower_map[k] for k, v in freq.most_common() if v < 2]
        themes = common + singles

        add_text_box(slide, Inches(0.5), y, Inches(12), Inches(0.35),
                     f"Q{qi+1}: {question}", font_size=Pt(11), colour=BLACK, bold=True)

        theme_text = " | ".join(themes[:8])
        add_text_box(slide, Inches(0.7), y + Inches(0.35), Inches(12), Inches(0.7),
                     theme_text, font_size=Pt(9), colour=DARK_GREY, bold=False)

        y += Inches(1.2)

    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


def build_back_cover_slide(prs, data):
    """Back cover — mirrors the front cover with William's contact details."""
    print("  Building Slide: Back Cover...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Solid green background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    # Logos — same as front cover, large and wide
    logo_height = Inches(1.8)
    logo_y = Inches(0.5)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(1.0), logo_y, height=logo_height)
    client_logo = os.path.join(os.path.dirname(os.path.abspath(__file__)), data.get("logo", "")) if data.get("logo") else None
    if client_logo and os.path.exists(client_logo):
        slide.shapes.add_picture(client_logo, Inches(10.5), logo_y, height=logo_height)

    # Title — same as front
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
                 "Strategic Readiness\nDiscussion Analysis",
                 font_size=Pt(40), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Company name
    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.8),
                 data["company"],
                 font_size=Pt(28), colour=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # Contact details
    add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.4),
                 "William Farmer",
                 font_size=Pt(16), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.4),
                 "Managing Director",
                 font_size=Pt(14), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
                 "0477 770 014",
                 font_size=Pt(14), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)


def build_how_to_read_slide(prs):
    """Insert a 'How to Read This Report' guide slide before individual results."""
    print("  Building Slide: How to Read This Report...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 "How to Read This Report", font_size=Pt(20), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    # ── Left column: Radar chart explanation ──
    left_x = Inches(0.6)
    y = Inches(1.3)

    add_text_box(slide, left_x, y, Inches(6), Inches(0.4),
                 "Understanding the Radar Charts", font_size=Pt(16), colour=BLACK, bold=True)
    y += Inches(0.55)

    explanations = [
        ("What it shows",
         "Each radar chart maps one person's scores across all competencies "
         "within a leadership profile. The scores are plotted as a shape on a spider web."),
        ("How to read the spokes",
         "Each spoke represents a competency (e.g. Self Belief, Motivation, Coaching). "
         "The centre of the chart is 0 and the outer edge is 10."),
        ("The coloured line",
         "This is the individual's score for each competency, connected to form a shape. "
         "A large, round shape indicates consistently strong performance. "
         "A spiky, uneven shape highlights clear strengths and development gaps."),
        ("The grey dashed line",
         "This represents the team average. Where the coloured line sits outside the grey, "
         "the person is above the team average. Where it sits inside, they are below."),
    ]

    for title, body in explanations:
        add_text_box(slide, left_x, y, Inches(5.8), Inches(0.3),
                     title, font_size=Pt(11), colour=DARK_GREY, bold=True)
        y += Inches(0.3)
        add_text_box(slide, left_x + Inches(0.1), y, Inches(5.7), Inches(0.65),
                     body, font_size=Pt(10), colour=BLACK, bold=False)
        y += Inches(0.65)

    # ── Right column: Colour key + score scale ──
    right_x = Inches(7.2)
    ry = Inches(1.3)

    add_text_box(slide, right_x, ry, Inches(5.5), Inches(0.4),
                 "Profile Colour Key", font_size=Pt(16), colour=BLACK, bold=True)
    ry += Inches(0.6)

    colour_key = [
        (GREEN, "Leadership DNA", "Foundational leadership competencies"),
        (YELLOW, "UYLP", "Interpersonal and communication skills"),
        (BLUE, "Leadership for Results", "Execution and performance competencies"),
    ]

    for colour, label, desc in colour_key:
        add_shape_rect(slide, right_x, ry + Pt(3), Inches(0.4), Inches(0.3), colour)
        add_text_box(slide, right_x + Inches(0.55), ry, Inches(4.5), Inches(0.3),
                     label, font_size=Pt(12), colour=BLACK, bold=True)
        ry += Inches(0.3)
        add_text_box(slide, right_x + Inches(0.55), ry, Inches(4.5), Inches(0.3),
                     desc, font_size=Pt(10), colour=DARK_GREY, bold=False)
        ry += Inches(0.5)

    # Score scale
    ry += Inches(0.3)
    add_text_box(slide, right_x, ry, Inches(5.5), Inches(0.4),
                 "Score Scale", font_size=Pt(16), colour=BLACK, bold=True)
    ry += Inches(0.55)

    score_key = [
        (GREEN, "9 — 10", "Strength"),
        (YELLOW, "7 — 8", "Competent"),
        (AMBER, "5 — 6", "Developing"),
        (RED, "1 — 4", "Development area"),
    ]

    for colour, range_text, meaning in score_key:
        add_shape_rect(slide, right_x, ry + Pt(2), Inches(0.4), Inches(0.25), colour)
        add_text_box(slide, right_x + Inches(0.55), ry, Inches(1.2), Inches(0.3),
                     range_text, font_size=Pt(11), colour=BLACK, bold=True)
        add_text_box(slide, right_x + Inches(1.8), ry, Inches(3), Inches(0.3),
                     meaning, font_size=Pt(11), colour=DARK_GREY, bold=False)
        ry += Inches(0.38)

    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


def generate_radar_chart(participant, team_avgs_by_profile, tmp_dir):
    """Generate a polished radar chart PNG — print-ready at 300dpi."""
    profile_configs = [
        {"key": "Leadership DNA", "color": "#84C44C", "label": "Leadership DNA"},
        {"key": "UYLP", "color": "#FFC708", "label": "UYLP"},
        {"key": "Leadership for Results", "color": "#0090CF", "label": "Leadership for Results"},
    ]

    fig, axes = plt.subplots(1, 3, figsize=(15, 5), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor('white')
    fig.subplots_adjust(wspace=0.55, top=0.82, bottom=0.1)

    for ax, pconf in zip(axes, profile_configs):
        profile_key = pconf["key"]
        scores = participant[profile_key]
        team_avgs = team_avgs_by_profile[profile_key]

        labels = list(scores.keys())
        short_labels = []
        for l in labels:
            if len(l) > 18:
                short_labels.append(l[:16] + '...')
            else:
                short_labels.append(l)

        values = list(scores.values())
        team_vals = [team_avgs.get(l, 0) for l in labels]
        n = len(labels)

        # Standard angles — NO rotation, so labels align consistently across charts
        angles = np.linspace(0, 2 * np.pi, n, endpoint=False).tolist()
        values_plot = values + [values[0]]
        team_plot = team_vals + [team_vals[0]]
        angles_plot = angles + [angles[0]]

        ax.set_ylim(0, 10)
        ax.set_yticks([2, 4, 6, 8, 10])
        ax.set_yticklabels(['2', '4', '6', '8', '10'], fontsize=5, color='#bbb')
        ax.set_xticks(angles)
        ax.set_xticklabels(short_labels, fontsize=5, color='#444')

        # Push labels outward for breathing room
        ax.tick_params(axis='x', pad=14)

        # Team average — subtle grey dashed fill
        ax.fill(angles_plot, team_plot, alpha=0.08, color='#888888')
        ax.plot(angles_plot, team_plot, color='#aaaaaa', linewidth=1, linestyle='--', label='Team Avg')

        # Individual — bold coloured fill
        ax.fill(angles_plot, values_plot, alpha=0.2, color=pconf["color"])
        ax.plot(angles_plot, values_plot, color=pconf["color"], linewidth=2.5, label=pconf["label"])

        ax.set_title(pconf["label"], fontsize=9, fontweight='bold', pad=22, color='#333')
        ax.legend(loc='lower center', bbox_to_anchor=(0.5, -0.2), fontsize=6,
                  framealpha=0.9, ncol=2)
        ax.grid(True, color='#e8e8e8', linewidth=0.4)
        ax.spines['polar'].set_color('#ddd')

    img_path = os.path.join(tmp_dir, f"radar_{participant['name'].replace(' ', '_')}.png")
    plt.savefig(img_path, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return img_path


def build_individual_slide(prs, data, participant, tmp_dir):
    """Build two slides per participant: summary + radar chart — print-ready."""
    name = participant["name"]
    print(f"  Building Individual Slide: {name}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar — taller, no crowding
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 name, font_size=Pt(22), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))

    # ── Left column: Strategic responses ──
    left_x = Inches(0.4)
    y = Inches(1.1)
    questions = data["questions"]
    q_keys = ["Q1", "Q2", "Q3", "Q4", "Q5"]

    add_text_box(slide, left_x, y, Inches(6), Inches(0.4),
                 "Strategic Readiness Responses", font_size=Pt(13), colour=BLACK, bold=True)
    y += Inches(0.45)

    for qi, (qkey, question) in enumerate(zip(q_keys, questions)):
        responses = participant["strategic"][qkey]
        add_text_box(slide, left_x, y, Inches(6.2), Inches(0.3),
                     f"Q{qi+1}: {question}", font_size=Pt(8.5), colour=DARK_GREY, bold=True)
        y += Inches(0.26)
        response_text = ", ".join(responses)
        add_text_box(slide, left_x + Inches(0.15), y, Inches(6), Inches(0.5),
                     response_text, font_size=Pt(8.5), colour=BLACK, bold=False)
        lines_est = max(1, len(response_text) // 75 + 1)
        y += Inches(0.22 * lines_est + 0.08)

    # ── Right column: Score summary ──
    right_x = Inches(7.2)
    ry = Inches(1.1)

    add_text_box(slide, right_x, ry, Inches(5.5), Inches(0.4),
                 "Score Summary", font_size=Pt(13), colour=BLACK, bold=True)
    ry += Inches(0.5)

    for profile in PROFILES:
        avg = compute_participant_profile_avg(participant, profile)
        colour = avg_score_colour(avg)
        add_shape_rect(slide, right_x, ry + Pt(2), Inches(0.5), Inches(0.3), colour)
        txt_colour = BLACK if 7 <= avg < 9 else WHITE
        add_text_box(slide, right_x, ry + Pt(2), Inches(0.5), Inches(0.3),
                     f"{avg}", font_size=Pt(11), colour=txt_colour, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, right_x + Inches(0.6), ry, Inches(4.5), Inches(0.35),
                     profile, font_size=Pt(11), colour=BLACK, bold=False)
        ry += Inches(0.4)

    ry += Inches(0.25)

    # Top 3 Strengths and Development Areas
    all_scores = []
    for profile in PROFILES:
        for comp, score in participant[profile].items():
            all_scores.append((comp, score, profile))

    sorted_high = sorted(all_scores, key=lambda x: x[1], reverse=True)
    sorted_low = sorted(all_scores, key=lambda x: x[1])

    seen = set()
    top3 = []
    for comp, score, profile in sorted_high:
        if comp not in seen and len(top3) < 3:
            top3.append((comp, score, profile))
            seen.add(comp)

    seen = set()
    bottom3 = []
    for comp, score, profile in sorted_low:
        if comp not in seen and len(bottom3) < 3:
            bottom3.append((comp, score, profile))
            seen.add(comp)

    add_text_box(slide, right_x, ry, Inches(5.5), Inches(0.35),
                 "Top 3 Strengths", font_size=Pt(12), colour=BLACK, bold=True)
    ry += Inches(0.35)
    for comp, score, profile in top3:
        add_shape_rect(slide, right_x, ry + Pt(3), Inches(0.12), Inches(0.22), GREEN)
        add_text_box(slide, right_x + Inches(0.2), ry, Inches(5.3), Inches(0.3),
                     f"{comp} ({score})", font_size=Pt(10), colour=BLACK, bold=False)
        ry += Inches(0.3)

    ry += Inches(0.15)
    add_text_box(slide, right_x, ry, Inches(5.5), Inches(0.35),
                 "Top 3 Development Areas", font_size=Pt(12), colour=BLACK, bold=True)
    ry += Inches(0.35)
    for comp, score, profile in bottom3:
        add_shape_rect(slide, right_x, ry + Pt(3), Inches(0.12), Inches(0.22), RED)
        add_text_box(slide, right_x + Inches(0.2), ry, Inches(5.3), Inches(0.3),
                     f"{comp} ({score})", font_size=Pt(10), colour=BLACK, bold=False)
        ry += Inches(0.3)

    add_shape_rect(slide, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)

    # ── Radar chart slide ──
    team_avgs_by_profile = {}
    for profile in PROFILES:
        team_avgs_by_profile[profile] = compute_team_averages(data, profile)

    radar_img = generate_radar_chart(participant, team_avgs_by_profile, tmp_dir)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape_rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(0.9), HEADER_BG)
    add_text_box(slide2, Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
                 f"{name} — Score Profile", font_size=Pt(18), colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide2.shapes.add_picture(LOGO_WHITE, Inches(11.5), Inches(0.15), height=Inches(0.55))

    slide2.shapes.add_picture(radar_img, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.6))

    add_shape_rect(slide2, Inches(0), Inches(7.3), Inches(13.333), Pt(3), HEADER_BG)


# ─── Main ────────────────────────────────────────────────────────────────────

def generate_srd_pptx(data):
    """Generate the full SRD analysis presentation.

    Slide order:
      1. Title (cover)
      2-17. Individual participant slides (summary + radar per person)
      18. Team Overview
      19-24. Score Tables + Bar Charts (Leadership DNA, UYLP, LFR)
      25. Collective Themes
    """
    company = data["company"]
    output_dir = os.path.expanduser("~/Desktop")
    output_path = os.path.join(output_dir, f"SRD Analysis - {company}.pptx")

    print(f"\nGenerating SRD Analysis for: {company}")
    print(f"Output: {output_path}\n")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    build_title_slide(prs, data)

    # Slide 2: Collective Themes
    build_collective_themes_slide(prs, data)

    # Slide 3: Team Overview (strengths + development areas)
    build_team_overview_slide(prs, data)

    # Slide 4: How to Read This Report
    build_how_to_read_slide(prs)

    # Slides 5+: Individual participant slides (people before detailed data)
    tmp_dir = tempfile.mkdtemp(prefix="srd_radar_")
    print(f"  Radar charts temp dir: {tmp_dir}")
    for participant in data["participants"]:
        build_individual_slide(prs, data, participant, tmp_dir)

    # Score tables + bar charts at the back
    build_score_table_slide(prs, data, "Leadership DNA")
    build_bar_chart_slide(prs, data, "Leadership DNA")
    build_score_table_slide(prs, data, "UYLP")
    build_bar_chart_slide(prs, data, "UYLP")
    build_score_table_slide(prs, data, "Leadership for Results")
    build_bar_chart_slide(prs, data, "Leadership for Results")

    # Back cover — mirrors the front with contact details
    build_back_cover_slide(prs, data)

    # Add slide numbers to all slides except the title (slide 1)
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        add_slide_number(slide, i + 1)

    prs.save(output_path)
    print(f"\nDone! {len(prs.slides)} slides generated.")
    print(f"File saved to: {output_path}")

    import shutil
    shutil.rmtree(tmp_dir, ignore_errors=True)

    return output_path


if __name__ == "__main__":
    company_name = sys.argv[1] if len(sys.argv) > 1 else None
    data = load_data(company_name)
    generate_srd_pptx(data)
