#!/usr/bin/env python3
"""
Generate Executive Coaching Report PowerPoint Presentation.

Generates a branded Dale Carnegie Queensland PPTX from exec-coaching-data.json.
New clients can be added as additional objects in the JSON array.

Usage:
    python3 generate-exec-coaching-pptx.py                  # Uses first client in JSON
    python3 generate-exec-coaching-pptx.py "Drew Morland"   # Specify client by name
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Brand Colours ───────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
YELLOW     = RGBColor(0xFF, 0xC7, 0x08)
GREEN      = RGBColor(0x84, 0xC4, 0x4C)
BLUE       = RGBColor(0x00, 0x90, 0xCF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0x99, 0x99, 0x99)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE     = RGBColor(0xE8, 0x86, 0x2A)
AMBER      = RGBColor(0xFF, 0xA5, 0x00)
RED        = RGBColor(0xFF, 0x63, 0x47)

# Header bar colour
HEADER_BG = GREEN

# Logos
LOGO_WHITE = os.path.expanduser("~/Dropbox/2026/Marketing/Logos & Badges/White DC Logo.png")
LOGO_BLACK = os.path.expanduser("~/Dropbox/2026/Marketing/Logos & Badges/Black DC Logo No Background.png")

# Slide size (widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Data Loading ────────────────────────────────────────────────────────────

def load_data(client_name=None):
    """Load client data from exec-coaching-data.json."""
    json_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "exec-coaching-data.json")
    with open(json_path) as f:
        all_data = json.load(f)
    clients = all_data["clients"]
    if client_name:
        for c in clients:
            if c["name"].lower() == client_name.lower():
                return c
        print(f"Client '{client_name}' not found. Available: {[c['name'] for c in clients]}")
        sys.exit(1)
    return clients[0]


def get_client_logo(data):
    """Resolve client logo path."""
    logo_path = data.get("client_logo", "")
    if logo_path:
        return os.path.expanduser(logo_path)
    return ""


def get_quarterly_plan(data):
    """Convert quarterly plan from JSON format to tuples."""
    result = []
    for q in data.get("quarterly_plan", []):
        result.append((q["label"], q["theme"], q["items"]))
    return result


# ─── Helper Functions ────────────────────────────────────────────────────────

def add_shape_rect(slide, left, top, width, height, fill_colour):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 colour=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = colour
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines, font_size=Pt(12),
                           colour=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                           line_spacing=Pt(18), bullet=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(4)
        if line_spacing:
            p.line_spacing = line_spacing
        prefix = "\u2022  " if bullet else ""
        run = p.add_run()
        run.text = prefix + line
        run.font.size = font_size
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = "Arial"
    return txBox


def add_slide_header(slide, title_text, font_size=Pt(18)):
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.9), HEADER_BG)
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(10.2), Inches(0.65),
                 title_text, font_size=font_size, colour=WHITE, bold=True)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(11.3), Inches(0.12), height=Inches(0.6))


def add_slide_footer(slide):
    add_shape_rect(slide, Inches(0), Inches(7.3), SLIDE_W, Pt(3), HEADER_BG)


def add_slide_number(slide, number):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
                 str(number), font_size=Pt(10), colour=GREY, bold=False, alignment=PP_ALIGN.RIGHT)


def score_bar_colour(score):
    if score >= 4.4:
        return GREEN
    elif score >= 4.0:
        return YELLOW
    else:
        return ORANGE


def disc_letter_colour(letter):
    colours = {"D": ORANGE, "I": YELLOW, "S": GREEN, "C": BLUE}
    return colours.get(letter, WHITE)


def add_client_logo(slide, data, logo_y):
    """Add client logo to a slide if available."""
    client_logo = get_client_logo(data)
    if client_logo and os.path.exists(client_logo):
        slide.shapes.add_picture(client_logo, Inches(9.5), logo_y + Inches(0.4), width=Inches(2.8))


# ─── Slide Builders ──────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    """Slide 1: Green cover with DC logo + client logo (same layout as SRD/Golden Cockerel)."""
    print("  Building Slide 1: Cover...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    # Logos — large, spread wide
    logo_height = Inches(1.8)
    logo_y = Inches(0.5)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(1.0), logo_y, height=logo_height)
    add_client_logo(slide, data, logo_y)

    # Main title
    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
                 "Executive Coaching\nReport",
                 font_size=Pt(40), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Name
    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.8),
                 data["name"],
                 font_size=Pt(28), colour=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    # Role
    add_text_box(slide, Inches(1.5), Inches(5.4), Inches(10.3), Inches(0.4),
                 data["role"],
                 font_size=Pt(18), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

    # Footer
    add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.4),
                 f"Coaching Session \u2014 {data['session']}",
                 font_size=Pt(13), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10.3), Inches(0.4),
                 "Prepared by William Farmer, Dale Carnegie Queensland",
                 font_size=Pt(13), colour=WHITE, bold=False, alignment=PP_ALIGN.CENTER)


def build_360_overview_slide(prs, data):
    """Slide 2: 360 Feedback — Category Scores (horizontal bars)."""
    print("  Building Slide 2: 360 Feedback Overview...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "360 Feedback Snapshot \u2014 Category Scores")

    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(10), Inches(0.4),
                 "Leadership Skills for Success \u2014 Category Scores (out of 5.0)",
                 font_size=Pt(16), colour=DARK_GREY, bold=True)

    y_start = Inches(1.7)
    bar_height = Inches(0.45)
    gap = Inches(0.55)
    bar_max_width = Inches(8.0)
    label_width = Inches(3.0)

    for i, (label, score) in enumerate(data["category_scores"]):
        y = y_start + gap * i
        add_text_box(slide, Inches(0.5), y, label_width, bar_height,
                     label, font_size=Pt(13), colour=DARK_GREY, bold=False,
                     alignment=PP_ALIGN.RIGHT)

        track_left = Inches(3.7)
        add_shape_rect(slide, track_left, y + Pt(4), bar_max_width, Inches(0.32), LIGHT_GREY)

        fill_width = int(bar_max_width * (score / 5.0))
        bar_colour = score_bar_colour(score)
        add_shape_rect(slide, track_left, y + Pt(4), fill_width, Inches(0.32), bar_colour)

        add_text_box(slide, track_left + fill_width + Inches(0.15), y + Pt(2),
                     Inches(0.6), bar_height,
                     f"{score:.2f}", font_size=Pt(12), colour=DARK_GREY, bold=True)

    add_slide_footer(slide)


def build_strengths_slide(prs, data):
    """Slide 3: Top Strengths from team feedback."""
    print("  Building Slide 3: Top Strengths...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Top Strengths \u2014 Team Feedback")

    y = Inches(1.3)
    for text, score in data["top_strengths"]:
        add_text_box(slide, Inches(0.8), y, Inches(10), Inches(0.5),
                     text, font_size=Pt(14), colour=DARK_GREY, bold=False)
        add_text_box(slide, Inches(11.0), y, Inches(1.5), Inches(0.5),
                     f"{score:.2f} / 5.00", font_size=Pt(14), colour=GREEN, bold=True,
                     alignment=PP_ALIGN.RIGHT)
        y += Inches(0.6)
        add_shape_rect(slide, Inches(0.8), y, Inches(11.5), Pt(1), LIGHT_GREY)
        y += Inches(0.2)

    first_name = data["name"].split()[0]
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(11), Inches(0.5),
                 f"{first_name}\u2019s team consistently rates him highest on trust, positive attitude, and genuine interest in people.",
                 font_size=Pt(13), colour=GREY, bold=False)

    add_slide_footer(slide)


def build_gaps_slide(prs, data):
    """Slide 4: Key Gaps — Self vs Team Discrepancies."""
    print("  Building Slide 4: Key Gaps...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Key Gaps \u2014 Self vs Team Discrepancies")

    first_name = data["name"].split()[0]
    add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
                 f"Items where {first_name} rated himself LOW but his team rated him HIGH:",
                 font_size=Pt(13), colour=GREY, bold=False)

    y = Inches(1.7)
    add_text_box(slide, Inches(0.5), y, Inches(5.5), Inches(0.35),
                 "Competency", font_size=Pt(11), colour=GREY, bold=True)
    add_text_box(slide, Inches(6.5), y, Inches(1.5), Inches(0.35),
                 "Self Score", font_size=Pt(11), colour=GREY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.2), y, Inches(1.5), Inches(0.35),
                 "Team Score", font_size=Pt(11), colour=GREY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(10.0), y, Inches(1.5), Inches(0.35),
                 "Gap", font_size=Pt(11), colour=GREY, bold=True, alignment=PP_ALIGN.CENTER)

    y += Inches(0.45)
    add_shape_rect(slide, Inches(0.5), y, Inches(11.5), Pt(1), DARK_GREY)
    y += Inches(0.15)

    for comp, self_score, team_score, gap in data["key_gaps"]:
        add_text_box(slide, Inches(0.5), y, Inches(5.5), Inches(0.4),
                     f"\u201c{comp}\u201d", font_size=Pt(12), colour=DARK_GREY, bold=False)
        add_text_box(slide, Inches(6.5), y, Inches(1.5), Inches(0.4),
                     f"{self_score:.1f}", font_size=Pt(12), colour=ORANGE, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(8.2), y, Inches(1.5), Inches(0.4),
                     f"{team_score:.2f}", font_size=Pt(12), colour=GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(10.0), y, Inches(1.5), Inches(0.4),
                     f"+{gap:.2f}", font_size=Pt(12), colour=GREEN, bold=True,
                     alignment=PP_ALIGN.CENTER)
        y += Inches(0.5)

    add_shape_rect(slide, Inches(0.5), y + Inches(0.2), Inches(11.5), Inches(0.7), LIGHT_GREY)
    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(11), Inches(0.5),
                 f"These positive gaps show {first_name} consistently underrates himself. His team sees significantly more capability than he self-perceives \u2014 a pattern consistent with his {data['disc_style']} DISC style.",
                 font_size=Pt(12), colour=DARK_GREY, bold=False)

    add_slide_footer(slide)


def build_attention_areas_slide(prs, data):
    """Slide 5: Areas Needing Attention (lowest team scores)."""
    print("  Building Slide 5: Areas Needing Attention...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Areas Needing Attention \u2014 Lowest Team Scores")

    y = Inches(1.5)
    bar_max_width = Inches(7.5)

    for label, score in data["areas_needing_attention"]:
        add_text_box(slide, Inches(0.5), y, Inches(4.5), Inches(0.45),
                     f"\u201c{label}\u201d", font_size=Pt(13), colour=DARK_GREY, bold=False,
                     alignment=PP_ALIGN.RIGHT)

        track_left = Inches(5.3)
        add_shape_rect(slide, track_left, y + Pt(6), bar_max_width, Inches(0.3), LIGHT_GREY)

        fill_width = int(bar_max_width * (score / 5.0))
        add_shape_rect(slide, track_left, y + Pt(6), fill_width, Inches(0.3), ORANGE)

        add_text_box(slide, track_left + fill_width + Inches(0.1), y + Pt(2),
                     Inches(0.8), Inches(0.4),
                     f"{score:.2f}", font_size=Pt(12), colour=ORANGE, bold=True)
        y += Inches(0.8)

    add_text_box(slide, Inches(0.8), y + Inches(0.3), Inches(11), Inches(0.5),
                 "Process skills, planning, and presentation delivery are the key development areas.",
                 font_size=Pt(13), colour=GREY, bold=False)

    add_slide_footer(slide)


def build_disc_profile_slide(prs, data):
    """Slide 6: DISC Profile Summary."""
    print("  Building Slide 6: DISC Profile...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "DISC Profile Summary")

    add_text_box(slide, Inches(0.5), Inches(1.2), Inches(2), Inches(0.6),
                 data["disc_style"], font_size=Pt(32), colour=YELLOW, bold=True)
    add_text_box(slide, Inches(2.7), Inches(1.35), Inches(4), Inches(0.4),
                 f'Behavioural Style: "{data["disc_label"]}"',
                 font_size=Pt(16), colour=DARK_GREY, bold=False)

    box_y = Inches(2.0)
    add_shape_rect(slide, Inches(0.5), box_y, Inches(5.8), Inches(1.8), LIGHT_GREY)
    add_text_box(slide, Inches(0.5), box_y + Inches(0.1), Inches(5.8), Inches(0.3),
                 "ADAPTED STYLE (WORK)", font_size=Pt(10), colour=GREY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    disc_x_start = Inches(1.0)
    for j, (letter, val) in enumerate(data["disc_adapted"].items()):
        x = disc_x_start + Inches(1.2) * j
        add_text_box(slide, x, box_y + Inches(0.5), Inches(1.0), Inches(0.5),
                     letter, font_size=Pt(24), colour=disc_letter_colour(letter),
                     bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, box_y + Inches(1.1), Inches(1.0), Inches(0.4),
                     str(val), font_size=Pt(16), colour=DARK_GREY,
                     bold=False, alignment=PP_ALIGN.CENTER)

    add_shape_rect(slide, Inches(7.0), box_y, Inches(5.8), Inches(1.8), LIGHT_GREY)
    add_text_box(slide, Inches(7.0), box_y + Inches(0.1), Inches(5.8), Inches(0.3),
                 "NATURAL STYLE", font_size=Pt(10), colour=GREY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    disc_x_start2 = Inches(7.5)
    for j, (letter, val) in enumerate(data["disc_natural"].items()):
        x = disc_x_start2 + Inches(1.2) * j
        add_text_box(slide, x, box_y + Inches(0.5), Inches(1.0), Inches(0.5),
                     letter, font_size=Pt(24), colour=disc_letter_colour(letter),
                     bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, box_y + Inches(1.1), Inches(1.0), Inches(0.4),
                     str(val), font_size=Pt(16), colour=DARK_GREY,
                     bold=False, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.1), Inches(5.5), Inches(0.4),
                 "Key Characteristics", font_size=Pt(15), colour=BLACK, bold=True)
    add_multiline_text_box(slide, Inches(0.5), Inches(4.5), Inches(5.5), Inches(2.5),
                           data["disc_characteristics"], font_size=Pt(12),
                           colour=DARK_GREY, bullet=True)

    add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.8), Inches(0.4),
                 "Focus Areas", font_size=Pt(15), colour=BLACK, bold=True)
    add_multiline_text_box(slide, Inches(7.0), Inches(4.5), Inches(5.8), Inches(2.5),
                           data["disc_focus_areas"][:5], font_size=Pt(11),
                           colour=DARK_GREY, bullet=True, line_spacing=Pt(16))

    add_slide_footer(slide)


def build_victories_slide(prs, data):
    """Slide 7: Coaching Journey — Victories."""
    print("  Building Slide 7: Victories...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Coaching Journey \u2014 Victories")

    add_shape_rect(slide, Inches(0.5), Inches(1.1), Inches(0.15), Inches(5.5), GREEN)

    add_multiline_text_box(slide, Inches(0.9), Inches(1.1), Inches(11.5), Inches(5.8),
                           data["victories"], font_size=Pt(13),
                           colour=DARK_GREY, bullet=True, line_spacing=Pt(20))

    add_slide_footer(slide)


def build_challenges_slide(prs, data):
    """Slide 8: Coaching Journey — Challenges Worked Through."""
    print("  Building Slide 8: Challenges...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Coaching Journey \u2014 Challenges Worked Through")

    add_shape_rect(slide, Inches(0.5), Inches(1.1), Inches(0.15), Inches(4.5), ORANGE)

    add_multiline_text_box(slide, Inches(0.9), Inches(1.1), Inches(11.5), Inches(5.0),
                           data["challenges"], font_size=Pt(13),
                           colour=DARK_GREY, bullet=True, line_spacing=Pt(20))

    add_slide_footer(slide)


def build_frameworks_slide(prs, data):
    """Slide 9: Key Frameworks Introduced."""
    print("  Building Slide 9: Frameworks...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Key Frameworks Introduced")

    add_shape_rect(slide, Inches(0.5), Inches(1.1), Inches(0.15), Inches(3.5), BLUE)

    add_multiline_text_box(slide, Inches(0.9), Inches(1.1), Inches(11.5), Inches(4.0),
                           data["frameworks"], font_size=Pt(14),
                           colour=DARK_GREY, bullet=True, line_spacing=Pt(24))

    first_name = data["name"].split()[0]
    add_text_box(slide, Inches(0.9), Inches(5.0), Inches(11), Inches(0.5),
                 f"These frameworks form the foundation of {first_name}\u2019s ongoing development and are revisited in each coaching session.",
                 font_size=Pt(13), colour=GREY, bold=False)

    add_slide_footer(slide)


def build_quarterly_plan_slide(prs, data, quarter_idx):
    """Build a slide for one quarter of the coaching plan."""
    plan = get_quarterly_plan(data)
    q_label, theme, items = plan[quarter_idx]
    slide_num = 10 + quarter_idx
    print(f"  Building Slide {slide_num}: {q_label}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, f"Coaching Plan \u2014 {q_label}")

    add_shape_rect(slide, Inches(0.5), Inches(1.2), Inches(2.2), Inches(0.5), YELLOW)
    add_text_box(slide, Inches(0.5), Inches(1.22), Inches(2.2), Inches(0.5),
                 q_label, font_size=Pt(14), colour=BLACK, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(3.0), Inches(1.25), Inches(6), Inches(0.45),
                 f'"{theme}"', font_size=Pt(18), colour=WHITE, bold=True)

    add_multiline_text_box(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
                           items, font_size=Pt(15),
                           colour=DARK_GREY, bullet=True, line_spacing=Pt(28))

    add_slide_footer(slide)


def build_comments_slide(prs, data):
    """Direct Report Comments from 360."""
    print("  Building Slide: Direct Report Comments...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Direct Report Comments \u2014 360 Feedback")

    y = Inches(1.3)
    for comment in data["direct_report_comments"]:
        add_shape_rect(slide, Inches(0.5), y, Pt(4), Inches(0.65), YELLOW)

        add_text_box(slide, Inches(0.8), y + Pt(4), Inches(11.5), Inches(0.6),
                     f"\u201c{comment}\u201d", font_size=Pt(12), colour=DARK_GREY,
                     bold=False)
        y += Inches(0.85)

    add_slide_footer(slide)


def build_back_cover_slide(prs, data):
    """Back cover — mirrors front cover with contact details."""
    print("  Building Slide: Back Cover...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = GREEN

    logo_height = Inches(1.8)
    logo_y = Inches(0.5)
    if os.path.exists(LOGO_WHITE):
        slide.shapes.add_picture(LOGO_WHITE, Inches(1.0), logo_y, height=logo_height)
    add_client_logo(slide, data, logo_y)

    add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(1.2),
                 "Executive Coaching\nReport",
                 font_size=Pt(40), colour=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10.3), Inches(0.8),
                 data["name"],
                 font_size=Pt(28), colour=BLACK, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.4),
                 "William Farmer", font_size=Pt(16), colour=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(6.1), Inches(10.3), Inches(0.4),
                 "Managing Director", font_size=Pt(14), colour=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
                 "0477 770 014", font_size=Pt(14), colour=WHITE, bold=False,
                 alignment=PP_ALIGN.CENTER)


# ─── Main ────────────────────────────────────────────────────────────────────

def generate_exec_coaching_pptx(data):
    """Generate the full executive coaching report presentation.

    Slide order:
      1.  Cover (green, DC + client logos)
      2.  360 Feedback — Category Scores
      3.  Top Strengths
      4.  Key Gaps — Self vs Team
      5.  Areas Needing Attention
      6.  DISC Profile Summary
      7.  Victories
      8.  Challenges Worked Through
      9.  Key Frameworks Introduced
      10-13. Coaching Plan (Q1-Q4)
      14. Direct Report Comments
      15. Back Cover
    """
    output_dir = os.path.expanduser("~/Desktop")
    output_path = os.path.join(output_dir, f"Executive Coaching Report - {data['name']}.pptx")

    print(f"\nGenerating Executive Coaching Report for: {data['name']}")
    print(f"Output: {output_path}\n")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, data)
    build_360_overview_slide(prs, data)
    build_strengths_slide(prs, data)
    build_gaps_slide(prs, data)
    build_attention_areas_slide(prs, data)
    build_disc_profile_slide(prs, data)
    build_victories_slide(prs, data)
    build_challenges_slide(prs, data)
    build_frameworks_slide(prs, data)

    plan = get_quarterly_plan(data)
    for i in range(len(plan)):
        build_quarterly_plan_slide(prs, data, i)

    build_comments_slide(prs, data)
    build_back_cover_slide(prs, data)

    # Add slide numbers (skip cover)
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        add_slide_number(slide, i + 1)

    prs.save(output_path)
    print(f"\nDone! {len(prs.slides)} slides generated.")
    print(f"File saved to: {output_path}")

    return output_path


if __name__ == "__main__":
    client_name = sys.argv[1] if len(sys.argv) > 1 else None
    data = load_data(client_name)
    generate_exec_coaching_pptx(data)
